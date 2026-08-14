#!/usr/bin/env python3
"""Build the 16-slide EditFlow evidence-first AgentRig GOAI 2026 review deck."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs/competition/AgentRig-GOAI-2026-初赛方案.pptx"
LIVE = ROOT / "docs/competition/assets/live"
OVERVIEW = LIVE / "editflow-01-overview.png"
CANDIDATE = LIVE / "editflow-02-candidate-run.png"
TIMELINE = LIVE / "editflow-03-timeline.png"
ACCEPTANCE = LIVE / "editflow-04-acceptance.png"
ASSISTANT = LIVE / "editflow-05-assistant.png"

W, H, TOTAL = 13.333, 7.5, 16
INK, TEXT, MUTED, FAINT = "121617", "2B312F", "67706B", "929A95"
DARK, PAPER, WHITE, LINE = "171B1F", "F7F6F2", "FFFFFF", "D8DDD9"
BLUE, BLUE_SOFT = "285CF5", "E9EFFF"
GREEN, GREEN_SOFT = "17745A", "E7F3EE"
AMBER, AMBER_SOFT = "A76509", "F8EEDB"
RED, RED_SOFT = "C84438", "F8E7E4"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def box(slide, x, y, w, h, *, fill=WHITE, stroke=LINE, width=0.7):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(stroke)
        shape.line.width = Pt(width)
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    *,
    size=15,
    color=TEXT,
    bold=False,
    font="PingFang SC",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    spacing=1.08,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = spacing
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def bullets(slide, values: Iterable[str], x, y, w, h, *, size=12, color=TEXT, gap=9):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for index, value in enumerate(values):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"•  {value}"
        p.font.name = "PingFang SC"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
    return shape


def rule(slide, x, y, w, *, color=LINE, height=0.012):
    return box(slide, x, y, w, height, fill=color, stroke=None)


def chip(slide, value, x, y, w, *, fill=BLUE_SOFT, color=BLUE):
    box(slide, x, y, w, 0.36, fill=fill, stroke=None)
    text(
        slide,
        value,
        x,
        y,
        w,
        0.36,
        size=8.4,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def footer(slide, index, *, dark=False):
    color = "737A77" if dark else FAINT
    text(slide, "AgentRig · GOAI 2026 · REVIEW", 0.72, 7.17, 3.5, 0.15, size=7, color=color)
    text(
        slide,
        f"{index:02d} / {TOTAL:02d}",
        11.65,
        7.17,
        0.96,
        0.15,
        size=7,
        color=color,
        font="Helvetica Neue",
        align=PP_ALIGN.RIGHT,
    )


def new_slide(prs, index, section, title_value, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    text(slide, section, 0.72, 0.34, 5.2, 0.2, size=8.5, color=BLUE, bold=True)
    text(slide, title_value, 0.72, 0.71, 11.4, 0.5, size=27, color=INK, bold=True)
    if subtitle:
        text(slide, subtitle, 0.74, 1.27, 11.5, 0.25, size=10.2, color=MUTED)
    text(
        slide,
        f"{index:02d}",
        12.05,
        0.37,
        0.55,
        0.2,
        size=9,
        color=FAINT,
        bold=True,
        font="Helvetica Neue",
        align=PP_ALIGN.RIGHT,
    )
    rule(slide, 0.72, 1.67, 11.9)
    footer(slide, index)
    return slide


def picture(slide, path: Path, x, y, w, h, *, border=True):
    if not path.is_file():
        raise FileNotFoundError(path)
    if border:
        box(slide, x - 0.02, y - 0.02, w + 0.04, h + 0.04, fill=WHITE, stroke=LINE)
    return slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )


def stat(slide, value, label, x, y, w, *, color=BLUE):
    text(slide, value, x, y, w, 0.45, size=26, color=color, bold=True, font="Helvetica Neue")
    text(slide, label, x, y + 0.53, w, 0.28, size=9.2, color=MUTED)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)

    # 01 — cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(DARK)
    box(slide, 0.72, 0.55, 0.42, 0.42, fill=BLUE, stroke=None)
    text(
        slide,
        "A",
        0.72,
        0.54,
        0.42,
        0.42,
        size=15,
        color=WHITE,
        bold=True,
        font="Helvetica Neue",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    text(slide, "AgentRig", 1.31, 0.66, 2.2, 0.24, size=10, color=WHITE, bold=True)
    text(
        slide,
        "GOAI 2026 · Agent Infra",
        9.45,
        0.66,
        3.15,
        0.2,
        size=8.5,
        color="A8B0AC",
        align=PP_ALIGN.RIGHT,
    )
    text(
        slide,
        "让 Agent 的每次变化\n都有可发布的证据",
        0.72,
        1.47,
        6.1,
        1.6,
        size=33,
        color=WHITE,
        bold=True,
        spacing=1.02,
    )
    text(
        slide,
        "真实推理 · 受控副作用 · 行为级回归",
        0.75,
        3.42,
        5.8,
        0.34,
        size=15.5,
        color="CAD0CD",
    )
    rule(slide, 0.75, 4.08, 5.55, color="363C39")
    text(
        slide,
        "Codex 改变 Agent，AgentRig 把改变固化成\nCase、Run、Timeline、Sample 与 Gate。",
        0.75,
        4.38,
        5.45,
        0.8,
        size=12.7,
        color="A8B0AC",
        spacing=1.35,
    )
    picture(slide, ACCEPTANCE, 7.03, 1.27, 5.55, 3.47)
    chip(slide, "EDITFLOW · 30 / 30 PASS", 7.03, 5.02, 2.3, fill=GREEN_SOFT, color=GREEN)
    text(
        slide,
        "公开脱敏 Agent · 真实 Agno + DeepSeek · 本地可复现",
        7.03,
        5.55,
        5.5,
        0.34,
        size=10.2,
        color="A8B0AC",
    )
    text(
        slide,
        "EVIDENCE BEFORE CONFIDENCE",
        0.75,
        5.78,
        3.2,
        0.2,
        size=8.5,
        color=BLUE,
        bold=True,
        font="Helvetica Neue",
    )
    footer(slide, 1, dark=True)

    # 02 — problem
    slide = new_slide(
        prs,
        2,
        "真实问题",
        "Prompt 改一行，Agent 行为可能变五层",
        "最终回答相似，不代表工具选择、参数、顺序和结果引用仍然正确。",
    )
    text(
        slide,
        "一次成功 Demo，\n不足以批准一次发布。",
        0.78,
        2.14,
        4.55,
        1.0,
        size=25,
        color=INK,
        bold=True,
    )
    text(
        slide,
        "Agent 的概率性与工具副作用，使传统单元测试和纯 Mock 都无法独立回答发布问题。",
        0.8,
        3.48,
        4.25,
        0.82,
        size=12.8,
        color=MUTED,
        spacing=1.3,
    )
    box(slide, 0.78, 5.0, 4.3, 0.76, fill=DARK, stroke=None)
    text(
        slide,
        "目标：真实决策，受控执行，可审计结论",
        1.0,
        5.22,
        3.86,
        0.3,
        size=12.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    risks = [
        ("01", "工具选择", "调用哪个工具，是否过度路由"),
        ("02", "参数传播", "人物保护、素材 ID、图片引用"),
        ("03", "调用顺序", "检查、修改、搜索、应用、裁剪"),
        ("04", "执行成本", "图片、支付、发信与生产写入"),
    ]
    for i, (number, name, detail) in enumerate(risks):
        y = 2.04 + i * 1.02
        chip(slide, number, 5.65, y, 0.58)
        text(slide, name, 6.48, y + 0.03, 1.45, 0.28, size=13, color=INK, bold=True)
        text(slide, detail, 7.95, y + 0.04, 4.2, 0.34, size=11, color=MUTED)
        rule(slide, 5.66, y + 0.7, 6.45)

    # 03 — boundary
    slide = new_slide(
        prs,
        3,
        "产品边界",
        "AgentRig 不替 Agent 回答，只治理评测边界",
        "模型仍真实推理；工具结果来源、事实记录与裁决责任被独立控制。",
    )
    stages = [
        ("TARGET", "真实被测 Agent", "模型推理\nHTTP / SSE\nSession 与工具选择", BLUE, BLUE_SOFT),
        (
            "CONTROL",
            "工具结果 Provider",
            "Fixture\n审核 Sample\nSimulator / Real Tool",
            AMBER,
            AMBER_SOFT,
        ),
        (
            "EVIDENCE",
            "事实与裁决",
            "Manifest\nCell / Attempt / Event\nRule / Judge / Human",
            GREEN,
            GREEN_SOFT,
        ),
    ]
    for i, (tag, name, detail, accent, soft) in enumerate(stages):
        x = 0.78 + i * 4.18
        box(slide, x, 2.05, 3.62, 3.66, fill=WHITE, stroke=LINE)
        chip(slide, tag, x + 0.28, 2.3, 1.08, fill=soft, color=accent)
        text(slide, name, x + 0.28, 2.92, 2.95, 0.36, size=18, color=INK, bold=True)
        text(slide, detail, x + 0.28, 3.58, 3.0, 1.35, size=13.5, color=MUTED, spacing=1.4)
        if i < 2:
            text(
                slide,
                "→",
                x + 3.72,
                3.62,
                0.4,
                0.4,
                size=22,
                color=BLUE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    box(slide, 2.18, 6.04, 8.98, 0.55, fill=DARK, stroke=None)
    text(
        slide,
        "保留决策真实性　｜　隔离工具副作用　｜　裁决不污染事实",
        2.4,
        6.18,
        8.55,
        0.26,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 04 — closed loop
    slide = new_slide(
        prs,
        4,
        "平台闭环",
        "Identity → Asset → Run → Evidence → Gate",
        "每一层都有稳定身份；历史 Run 不被修复动作覆盖。",
    )
    flow = [
        ("01", "Identity", "Target\nPrompt SHA\nTool schema", BLUE, BLUE_SOFT),
        ("02", "Assets", "Case\nProfile\nSample", AMBER, AMBER_SOFT),
        ("03", "Manifest", "Canonical\nPreview\nHash", BLUE, BLUE_SOFT),
        ("04", "Execution", "Cell\nAttempt\nEvent", GREEN, GREEN_SOFT),
        ("05", "Decision", "Rule\nJudge\nHuman", RED, RED_SOFT),
        ("06", "Release", "Report\nGate\nAudit", GREEN, GREEN_SOFT),
    ]
    for i, (number, name, detail, accent, soft) in enumerate(flow):
        x = 0.77 + i * 2.02
        box(slide, x, 2.17, 1.72, 3.12, fill=WHITE, stroke=LINE)
        chip(slide, number, x + 0.24, 2.43, 0.5, fill=soft, color=accent)
        text(slide, name, x + 0.24, 3.06, 1.3, 0.3, size=13.5, color=INK, bold=True)
        text(slide, detail, x + 0.24, 3.72, 1.28, 1.05, size=11, color=MUTED, spacing=1.35)
        if i < 5:
            text(
                slide,
                "→",
                x + 1.73,
                3.55,
                0.3,
                0.3,
                size=16,
                color=BLUE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    text(
        slide,
        "同一事实链支持本地调试、团队评审、CI Gate 与生产回灌",
        2.15,
        5.84,
        9.0,
        0.35,
        size=14,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    text(
        slide,
        "不是一张绿色分数卡，而是一条可定位、可复跑、可追责的发布证据链。",
        2.0,
        6.35,
        9.35,
        0.28,
        size=10.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    # 05 — two entries
    slide = new_slide(
        prs,
        5,
        "双入口",
        "不同能力可以给出不同方案，共享同一种证据合同",
        "Codex 面向开发者闭环；Web 助手面向产品、运营与测试人员。",
    )
    box(slide, 0.75, 2.0, 5.95, 4.25, fill=DARK, stroke=None)
    text(slide, "DEVELOPER", 1.05, 2.28, 1.5, 0.22, size=8.5, color=BLUE, bold=True)
    text(slide, "Codex + 项目 Skill + MCP", 1.05, 2.75, 5.1, 0.4, size=20, color=WHITE, bold=True)
    bullets(
        slide,
        [
            "读取代码、Prompt Diff 和业务不变量",
            "查询、创建并治理 Case / Sample",
            "Preview 后按 Manifest 提交回归矩阵",
            "基于 Timeline 修复并输出验收结论",
        ],
        1.05,
        3.47,
        5.15,
        1.75,
        size=11.3,
        color="CFD4D1",
        gap=8,
    )
    chip(slide, "主入口 · 改软件 + 治回归", 1.05, 5.57, 2.35, fill=GREEN_SOFT, color=GREEN)
    box(slide, 6.98, 2.0, 5.6, 4.25, fill=WHITE, stroke=LINE)
    text(slide, "PRODUCT / QA", 7.28, 2.28, 1.8, 0.22, size=8.5, color=GREEN, bold=True)
    text(slide, "AgentRig 智能评测助手", 7.28, 2.75, 4.8, 0.4, size=20, color=INK, bold=True)
    bullets(
        slide,
        [
            "自然语言描述评测目标和成本边界",
            "模型生成自己的可编辑计划",
            "高成本执行必须人工确认",
            "复用 Target、Case、Run 与报告资产",
        ],
        7.28,
        3.47,
        4.72,
        1.75,
        size=11.3,
        color=MUTED,
        gap=8,
    )
    chip(slide, "辅助入口 · 无需 MCP 知识", 7.28, 5.57, 2.35, fill=AMBER_SOFT, color=AMBER)
    text(
        slide,
        "不要求两种模型生成同一方案；要求每个结论都能回到冻结资产和运行证据。",
        2.1,
        6.53,
        9.2,
        0.27,
        size=10.5,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 06 — EditFlow
    slide = new_slide(
        prs,
        6,
        "公开被测对象",
        "EditFlow：可复现、可脱敏、仍然足够真实",
        "以 lassist 的真实修图链路为基础，去除私有素材、版本与基础设施依赖。",
    )
    box(slide, 0.78, 2.02, 4.2, 4.5, fill=DARK, stroke=None)
    chip(slide, "OPEN DEMO AGENT", 1.07, 2.32, 1.65, fill="1D2B34", color="80A9FF")
    text(slide, "Agno + DeepSeek", 1.07, 2.94, 3.5, 0.4, size=21, color=WHITE, bold=True)
    text(
        slide,
        "PostgreSQL Session\nHTTP/SSE Target\nExternal tool execution\nLocal deterministic MCP",
        1.07,
        3.66,
        3.3,
        1.62,
        size=12.5,
        color="CFD4D1",
        spacing=1.45,
    )
    text(
        slide,
        "不上传图片 · 不调用第三方图片 API",
        1.07,
        5.83,
        3.45,
        0.28,
        size=10,
        color="8EA5FF",
        bold=True,
    )
    tools_list = [
        ("inspect_image", "读取图片属性与主体"),
        ("retouch_photo", "亮度、对比度等像素调整"),
        ("search_assets", "搜索受控素材并返回真实 ID"),
        ("apply_asset", "消费素材 ID 与最新图片引用"),
        ("crop_photo", "按目标比例裁剪"),
    ]
    for i, (name, detail) in enumerate(tools_list):
        y = 2.03 + i * 0.84
        chip(slide, f"0{i + 1}", 5.42, y, 0.48)
        text(slide, name, 6.18, y + 0.02, 2.15, 0.27, size=11.5, color=INK, bold=True, font="Menlo")
        text(slide, detail, 8.48, y + 0.03, 3.7, 0.3, size=10.5, color=MUTED)
        rule(slide, 5.42, y + 0.6, 6.75)
    box(slide, 5.42, 6.25, 6.75, 0.42, fill=GREEN_SOFT, stroke=None)
    text(
        slide,
        "真实：模型、协议、Session、工具决策　｜　受控：图片工具结果",
        5.58,
        6.35,
        6.42,
        0.22,
        size=9.5,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 07 — Before
    slide = new_slide(
        prs,
        7,
        "真实回归 · Before",
        "同一冻结 Case 跑 5 次：2 Pass / 3 Behavior Fail",
        "单次成功掩盖模型方差；父 Run completed 不等于业务验收通过。",
    )
    box(slide, 0.78, 2.02, 4.25, 4.72, fill=DARK, stroke=None)
    text(slide, "HEADLINE TASK", 1.08, 2.31, 1.7, 0.22, size=8.5, color=BLUE, bold=True)
    text(
        slide,
        "一步完成调亮、\n雪山背景与 4:5 裁剪，\n并保持人物。",
        1.08,
        2.9,
        3.4,
        1.48,
        size=20,
        color=WHITE,
        bold=True,
        spacing=1.22,
    )
    text(
        slide,
        "case_editflow_one_step_mixed_chain",
        1.08,
        5.06,
        3.4,
        0.42,
        size=9.3,
        color="AAB3AF",
        font="Menlo",
    )
    text(
        slide,
        "Manifest  bb5008e6…6766aa",
        1.08,
        5.68,
        3.25,
        0.24,
        size=9.3,
        color="AAB3AF",
        font="Menlo",
    )
    text(
        slide,
        "Prompt SHA 43157c…9572",
        1.08,
        6.08,
        3.25,
        0.24,
        size=9.3,
        color="AAB3AF",
        font="Menlo",
    )
    stat(slide, "2", "Pass", 5.6, 2.22, 1.2, color=GREEN)
    stat(slide, "3", "Behavior Fail", 7.22, 2.22, 1.7, color=RED)
    stat(slide, "5", "Independent Attempts", 9.5, 2.22, 2.3, color=BLUE)
    rule(slide, 5.58, 3.32, 6.5)
    text(slide, "失败轨迹", 5.6, 3.72, 1.3, 0.28, size=13.5, color=RED, bold=True)
    box(slide, 5.6, 4.18, 6.48, 1.2, fill=RED_SOFT, stroke=None)
    text(
        slide,
        "retouch_photo  →  inspect_image  →  search_assets  →  apply_asset  →  crop_photo",
        5.88,
        4.48,
        5.9,
        0.48,
        size=11.2,
        color=RED,
        bold=True,
        font="Menlo",
        align=PP_ALIGN.CENTER,
    )
    text(
        slide,
        "检查有时落在调亮之后，违反“先检查再修改”的冻结行为契约。",
        5.6,
        5.71,
        6.45,
        0.4,
        size=12,
        color=TEXT,
        bold=True,
    )
    text(
        slide, "Run  run_61e76493…8d9c03", 5.6, 6.28, 4.8, 0.25, size=9.4, color=MUTED, font="Menlo"
    )

    # 08 — case governance
    slide = new_slide(
        prs,
        8,
        "Codex + Skill",
        "不是盲目增加用例，而是按 Prompt Diff 治理风险",
        "项目 Skill 先冻结身份与不变量，再查询 AgentRig 资产并分类。",
    )
    classes = [
        ("REUSE", "复用", "已有用例已覆盖，直接进入矩阵", GREEN, GREEN_SOFT),
        ("CHANGE", "增强", "行为契约变化，增强现有断言", AMBER, AMBER_SOFT),
        ("NEW", "新增", "变更引入新风险，现场创建", BLUE, BLUE_SOFT),
        ("EXCLUDE", "排除", "与本次 Prompt Diff 无关", MUTED, "ECEFEC"),
    ]
    for i, (tag, name, detail, accent, soft) in enumerate(classes):
        x = 0.78 + i * 3.0
        box(slide, x, 2.05, 2.68, 2.15, fill=WHITE, stroke=LINE)
        chip(slide, tag, x + 0.25, 2.3, 0.92, fill=soft, color=accent)
        text(slide, name, x + 0.25, 2.93, 1.2, 0.3, size=16, color=INK, bold=True)
        text(slide, detail, x + 0.25, 3.46, 2.2, 0.5, size=10.2, color=MUTED, spacing=1.25)
    text(slide, "本次防回归矩阵", 0.8, 4.71, 2.2, 0.3, size=15, color=INK, bold=True)
    cases = [
        ("Change", "brighten_only_no_over_routing", "简单调亮不能拆成素材或裁剪"),
        ("Change", "asset_search_miss_no_invention", "搜索为空不能编造 asset_id"),
        ("Change", "preserve_subject_propagation", "人物保护必须跨工具传播"),
        ("New", "one_step_mixed_chain_boundary", "一步完成也不能绕过专用工具边界"),
    ]
    for i, (kind, name, detail) in enumerate(cases):
        y = 5.18 + i * 0.4
        text(
            slide,
            kind,
            0.82,
            y,
            0.78,
            0.22,
            size=8.6,
            color=BLUE if kind == "New" else AMBER,
            bold=True,
        )
        text(slide, name, 1.7, y, 3.55, 0.22, size=9.6, color=TEXT, font="Menlo", bold=True)
        text(slide, detail, 5.45, y, 6.45, 0.22, size=9.6, color=MUTED)
    text(
        slide,
        "Draft → 人工审核 → approved；创建者不能批准自己创建的资产。",
        2.1,
        6.77,
        9.1,
        0.22,
        size=9.8,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 09 — prompt identity
    slide = new_slide(
        prs,
        9,
        "最小修复",
        "只改模型可见边界，并用 Prompt SHA 证明 Candidate 已生效",
        "HTTP/SSE 适配层没有加入关键词路由，也没有为视频写死答案。",
    )
    box(slide, 0.78, 2.03, 7.25, 4.5, fill="0E1418", stroke=None)
    chip(slide, "PROMPT DIFF", 1.06, 2.3, 1.2, fill="1D2B34", color="80A9FF")
    diff = [
        "- retouch_photo 可完成复合图片编辑",
        "+ retouch_photo 仅处理亮度、对比度等像素调整",
        "+ 背景/素材必须 search_assets → apply_asset",
        "+ 比例必须 crop_photo",
        "+ 下一工具消费最新 output_image_ref",
        "+ preserve_subject 传播到相关工具",
    ]
    text(
        slide,
        "\n".join(diff),
        1.06,
        2.94,
        6.5,
        2.7,
        size=11.2,
        color="D7E0E4",
        font="Menlo",
        spacing=1.3,
    )
    text(
        slide,
        "修改范围：system.md + tool description",
        1.06,
        5.93,
        5.8,
        0.27,
        size=9.5,
        color="8EA5FF",
        font="Menlo",
    )
    text(slide, "BEFORE", 8.5, 2.15, 1.0, 0.22, size=8.5, color=RED, bold=True)
    text(slide, "43157c…9572", 8.5, 2.63, 3.42, 0.35, size=16, color=RED, bold=True, font="Menlo")
    rule(slide, 8.5, 3.28, 3.55)
    text(slide, "CANDIDATE", 8.5, 3.68, 1.35, 0.22, size=8.5, color=GREEN, bold=True)
    text(slide, "71adb3…f7fd", 8.5, 4.16, 3.42, 0.35, size=16, color=GREEN, bold=True, font="Menlo")
    box(slide, 8.5, 5.03, 3.55, 0.78, fill=GREEN_SOFT, stroke=None)
    text(
        slide,
        "SHA 不变，Candidate\n就没有验收资格",
        8.72,
        5.22,
        3.1,
        0.42,
        size=11,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        spacing=1.2,
    )
    text(
        slide,
        "仓库门禁：34 tests passed",
        8.5,
        6.21,
        3.55,
        0.24,
        size=10,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 10 — candidate matrix
    slide = new_slide(
        prs,
        10,
        "Candidate 回归",
        "同 Case 5/5；六 Case 矩阵 30/30",
        "每次 Attempt 都独立执行；没有用一个绿色总分掩盖方差。",
    )
    picture(slide, CANDIDATE, 0.78, 1.98, 9.15, 5.04)
    stat(slide, "6", "Cells", 10.25, 2.14, 0.9, color=BLUE)
    stat(slide, "30", "Attempts", 11.28, 2.14, 1.1, color=BLUE)
    rule(slide, 10.23, 3.2, 2.18)
    stat(slide, "30", "Pass", 10.25, 3.57, 1.0, color=GREEN)
    stat(slide, "0", "Fail", 11.28, 3.57, 1.0, color=RED)
    box(slide, 10.23, 4.73, 2.18, 1.1, fill=GREEN_SOFT, stroke=None)
    text(
        slide,
        "headline\n2 / 5 → 5 / 5",
        10.42,
        4.96,
        1.8,
        0.58,
        size=12.5,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        spacing=1.2,
    )
    text(
        slide,
        "run_3afacf…4793",
        10.23,
        6.18,
        2.18,
        0.24,
        size=8.3,
        color=MUTED,
        font="Menlo",
        align=PP_ALIGN.CENTER,
    )

    # 11 — timeline
    slide = new_slide(
        prs,
        11,
        "行为证据",
        "不只验证最终文字：完整 Timeline 可以下钻",
        "工具参数、Provider 来源、结果引用与裁决都关联同一 Attempt。",
    )
    picture(slide, TIMELINE, 0.78, 1.98, 8.95, 5.04)
    expected = [
        "inspect_image",
        "search_assets",
        "retouch_photo",
        "apply_asset",
        "crop_photo",
    ]
    for i, value in enumerate(expected):
        y = 2.08 + i * 0.68
        chip(slide, f"{i + 1}", 10.05, y, 0.44, fill=BLUE_SOFT, color=BLUE)
        text(
            slide, value, 10.72, y + 0.05, 1.65, 0.24, size=9.5, color=INK, bold=True, font="Menlo"
        )
    rule(slide, 10.05, 5.58, 2.28)
    text(
        slide,
        "✓ asset_id 来自搜索\n✓ preserve_subject 传播\n✓ image_ref 逐步更新",
        10.05,
        5.88,
        2.3,
        0.8,
        size=9.5,
        color=GREEN,
        bold=True,
        spacing=1.3,
    )

    # 12 — sample
    slide = new_slide(
        prs,
        12,
        "低副作用核心",
        "真实工具结果，如何变成可审核的零真实调用回归资产",
        "真实来源、人工责任和重复成本被放进同一条证据链。",
    )
    steps = [
        ("01", "Capture", "真实 MCP 调用 1 次", "run_db173d…59ef", BLUE, BLUE_SOFT),
        ("02", "Event", "持久化 real_tool 事实", "evt_32fb37…f80d", AMBER, AMBER_SOFT),
        ("03", "Sample", "Codex 创建 Draft", "sample_editflow…0814", BLUE, BLUE_SOFT),
        ("04", "Review", "人类批准来源与边界", "approved", RED, RED_SOFT),
        ("05", "Replay", "Sample-only × 5", "run_1a209c…0b90", GREEN, GREEN_SOFT),
    ]
    for i, (number, name, detail, ident, accent, soft) in enumerate(steps):
        x = 0.78 + i * 2.42
        box(slide, x, 2.1, 2.12, 3.55, fill=WHITE, stroke=LINE)
        chip(slide, number, x + 0.24, 2.38, 0.48, fill=soft, color=accent)
        text(slide, name, x + 0.24, 3.03, 1.55, 0.3, size=15, color=INK, bold=True)
        text(slide, detail, x + 0.24, 3.67, 1.62, 0.62, size=10.2, color=MUTED, spacing=1.25)
        text(
            slide, ident, x + 0.24, 4.69, 1.62, 0.5, size=8.2, color=accent, bold=True, font="Menlo"
        )
        if i < 4:
            text(
                slide,
                "→",
                x + 2.11,
                3.63,
                0.3,
                0.3,
                size=16,
                color=BLUE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    box(slide, 2.02, 6.03, 9.3, 0.62, fill=DARK, stroke=None)
    text(
        slide,
        "Replay：5 / 5 Sample hit　｜　该 Run 中 real_tool Provider Attempt = 0",
        2.24,
        6.2,
        8.86,
        0.28,
        size=11.2,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 13 — assistant
    slide = new_slide(
        prs,
        13,
        "普通用户入口",
        "自然语言规划，确认与提交分离",
        "无需理解 MCP；成本边界和人工责任仍由 Core 强制执行。",
    )
    picture(slide, ASSISTANT, 0.78, 1.98, 9.15, 5.04)
    stat(slide, "1", "Case", 10.25, 2.1, 0.9, color=BLUE)
    stat(slide, "3", "Attempts", 11.3, 2.1, 1.1, color=BLUE)
    rule(slide, 10.23, 3.16, 2.18)
    text(slide, "01  生成 Draft Plan", 10.25, 3.53, 2.0, 0.25, size=10.3, color=TEXT, bold=True)
    text(slide, "02  重复执行触发确认", 10.25, 4.02, 2.0, 0.25, size=10.3, color=TEXT, bold=True)
    text(slide, "03  confirm 不创建 Run", 10.25, 4.51, 2.0, 0.25, size=10.3, color=TEXT, bold=True)
    text(slide, "04  submit 才执行", 10.25, 5.0, 2.0, 0.25, size=10.3, color=TEXT, bold=True)
    box(slide, 10.23, 5.58, 2.18, 0.72, fill=GREEN_SOFT, stroke=None)
    text(
        slide,
        "run_a1c88b…23fb\n3 / 3 PASS",
        10.38,
        5.72,
        1.88,
        0.42,
        size=9.8,
        color=GREEN,
        bold=True,
        font="Menlo",
        align=PP_ALIGN.CENTER,
        spacing=1.2,
    )

    # 14 — platform
    slide = new_slide(
        prs,
        14,
        "平台优势",
        "不是 EditFlow 专用脚本，而是可替换的评测控制面",
        "框架、结果来源、裁决方式和消费入口彼此解耦。",
    )
    layers = [
        ("入口", "Codex / MCP · Web Assistant · CLI · CI", BLUE, BLUE_SOFT),
        ("Driver", "HTTP/SSE · OpenAI-compatible · AgentScope · AG-UI", AMBER, AMBER_SOFT),
        ("Provider", "Fixture · Approved Sample · Simulator · Real Tool", GREEN, GREEN_SOFT),
        ("Evidence", "Manifest · Cell · Attempt · Event · Timeline", BLUE, BLUE_SOFT),
        ("Decision", "Rule · Judge · Human Review · Report · Release Gate", RED, RED_SOFT),
    ]
    for i, (name, detail, accent, soft) in enumerate(layers):
        y = 2.02 + i * 0.83
        chip(slide, name, 0.82, y, 1.02, fill=soft, color=accent)
        box(slide, 2.08, y, 10.02, 0.56, fill=WHITE, stroke=LINE)
        text(slide, detail, 2.37, y + 0.14, 9.45, 0.27, size=11.5, color=TEXT, bold=i >= 3)
    chips = ["Project isolation", "Audit log", "Durable job", "OTLP → Case", "Failure patterns"]
    for i, value in enumerate(chips):
        chip(slide, value, 0.82 + i * 2.38, 6.35, 2.1, fill="ECEFEC", color=MUTED)
    text(
        slide,
        "核心模块可单独使用；AgentTeams 只在动态策展或语义裁决需要时启用。",
        2.1,
        6.84,
        9.1,
        0.22,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    # 15 — ecosystem
    slide = new_slide(
        prs,
        15,
        "开源生态",
        "把 AgentScope 的真实闭环经验，抽成可本地部署的开放合同",
        "AgentRig 不是复制一个内部页面，而是开放身份、资产、运行、证据和 Gate。",
    )
    columns = [
        (
            "AgentScope 实践",
            ["真实业务评测闭环", "工作台交互经验", "AgentTeams 协作范式"],
            BLUE,
            BLUE_SOFT,
        ),
        (
            "AgentRig 开源内核",
            ["协议无关 Driver", "低副作用 Provider", "Canonical Evidence"],
            GREEN,
            GREEN_SOFT,
        ),
        (
            "生态消费方式",
            ["Codex + 项目 Skill", "普通用户 Web Assistant", "CI / 发布 Gate"],
            AMBER,
            AMBER_SOFT,
        ),
    ]
    for i, (title_value, values, accent, soft) in enumerate(columns):
        x = 0.78 + i * 4.13
        box(slide, x, 2.04, 3.72, 3.9, fill=WHITE, stroke=LINE)
        chip(slide, f"0{i + 1}", x + 0.3, 2.33, 0.5, fill=soft, color=accent)
        text(slide, title_value, x + 0.3, 2.96, 3.05, 0.35, size=17, color=INK, bold=True)
        bullets(slide, values, x + 0.3, 3.66, 3.02, 1.44, size=11, color=MUTED, gap=10)
    box(slide, 1.25, 6.27, 10.83, 0.52, fill=DARK, stroke=None)
    text(
        slide,
        "公开 EditFlow 证明可复现　｜　真实 lassist 兼容验证证明不是只适配 Demo Agent",
        1.5,
        6.41,
        10.35,
        0.24,
        size=10.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 16 — ledger
    slide = new_slide(
        prs,
        16,
        "证据台账",
        "正式录制已闭环；台账中的每个 Run、Sample 与事件均来自干净录制库",
        "已证明的事实与尚未证明的边界必须同时出现在结论里。",
    )
    ledger = [
        ("Before", "2 / 5", "3 次行为失败", RED),
        ("Headline", "5 / 5", "同 Case + Manifest", GREEN),
        ("Matrix", "30 / 30", "6 Cells", GREEN),
        ("Sample replay", "5 / 5", "0 real_tool", GREEN),
        ("Web Assistant", "3 / 3", "确认与提交分离", GREEN),
        ("EditFlow", "34 tests", "non-live", BLUE),
    ]
    for i, (name, value, detail, accent) in enumerate(ledger):
        col, row = i % 3, i // 3
        x, y = 0.78 + col * 4.03, 2.02 + row * 1.15
        box(slide, x, y, 3.68, 0.92, fill=WHITE, stroke=LINE)
        text(slide, name, x + 0.24, y + 0.16, 1.35, 0.22, size=9.5, color=MUTED, bold=True)
        text(
            slide,
            value,
            x + 1.48,
            y + 0.13,
            1.05,
            0.3,
            size=14,
            color=accent,
            bold=True,
            font="Helvetica Neue",
        )
        text(slide, detail, x + 2.45, y + 0.17, 1.0, 0.28, size=8.6, color=MUTED)
    box(slide, 0.78, 4.72, 11.73, 1.02, fill=AMBER_SOFT, stroke=None)
    text(slide, "诚实边界", 1.04, 4.98, 1.18, 0.25, size=10, color=AMBER, bold=True)
    text(
        slide,
        "Sample 仅覆盖 inspect_image；动态引用等值以 Timeline 展示；正式录制必须生成新 ID；\n"
        "AgentTeams 外部 Live、目标容量与生产 SLO 不冒充本次 EditFlow 结论。",
        2.35,
        4.91,
        9.5,
        0.58,
        size=10.2,
        color=TEXT,
        bold=True,
        spacing=1.25,
    )
    box(slide, 1.52, 6.05, 10.3, 0.7, fill=DARK, stroke=None)
    text(
        slide,
        "Codex 负责思考和改变软件；AgentRig 负责让每次改变可复用、可执行、可追溯。",
        1.78,
        6.25,
        9.78,
        0.3,
        size=12,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    if len(prs.slides) != TOTAL:
        raise RuntimeError(f"Expected {TOTAL} slides, got {len(prs.slides)}")
    return prs


def main() -> None:
    for asset in (OVERVIEW, CANDIDATE, TIMELINE, ACCEPTANCE, ASSISTANT):
        if not asset.is_file():
            raise FileNotFoundError(asset)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    build().save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
